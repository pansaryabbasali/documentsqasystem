"""PPTX extraction via python-pptx, one TextBlock per slide.

Text frames and table cells are both captured — business decks put their
load-bearing numbers in tables. Speaker notes are deliberately skipped:
they're author-facing, not part of the presented document.
"""

from __future__ import annotations

from collections.abc import Iterator
from pathlib import Path

from pptx import Presentation

from .base import TextBlock


class PptxLoader:
    suffixes: tuple[str, ...] = (".pptx",)

    def load(self, path: Path) -> Iterator[TextBlock]:
        for number, slide in enumerate(Presentation(str(path)).slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text.strip())
                elif getattr(shape, "has_table", False):
                    parts.extend(
                        " | ".join(cell.text.strip() for cell in row.cells)
                        for row in shape.table.rows
                    )
            if parts:
                yield TextBlock(text="\n".join(parts), source=path.name, locator=f"slide {number}")
